import streamlit as st
import pandas as pd
import base64
import io
import numpy as np
import re
from PIL import Image
import matplotlib.pyplot as plt
# import spacy
import logging
import warnings
from nltk.corpus import stopwords
import nltk
import os
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Border, Side, Alignment, Font,PatternFill
from openpyxl.utils.dataframe import dataframe_to_rows # Add these imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches


# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Data preprocessing function (You can include your data preprocessing here)

# Function to create separate Excel sheets by Entity
def create_entity_sheets(data, writer):
    # Define a format with text wrap
    wrap_format = writer.book.add_format({'text_wrap': True})

    for Entity in data['Entity'].unique():
        entity_df = data[data['Entity'] == Entity]
        entity_df.to_excel(writer, sheet_name=Entity, index=False)
        worksheet = writer.sheets[Entity]
        worksheet.set_column(1, 4, 48, cell_format=wrap_format)
        # Calculate column widths based on the maximum content length in each column except columns 1 to 4
        max_col_widths = [
            max(len(str(value)) for value in entity_df[column])
            for column in entity_df.columns[5:]  # Exclude columns 1 to 4
        ]

        # Set the column widths dynamically for columns 5 onwards
        for col_num, max_width in enumerate(max_col_widths):
            worksheet.set_column(col_num + 5, col_num + 5, max_width + 2)  # Adding extra padding for readability       


def add_entity_info(ws, entity_info, start_row):
    for i, line in enumerate(entity_info.split('\n'), start=1):
        cell = ws.cell(row=start_row + i - 1, column=1)
        cell.value = line
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))
#         cell.alignment = Alignment(horizontal='center')  # Merge and center for all lines
#         ws.merge_cells(start_row=start_row + i - 1, start_column=1, end_row=start_row + i, end_column=5)
        
        # Apply specific formatting for Source line
        if line.startswith('Source:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=3, start_column=1, end_row=3, end_column=5)
            cell.font = Font(color="000000",name="Gill Sans")
            
        # Apply specific formatting for Source line
        if line.startswith('Entity:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans", bold=True )
            cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
            
        # Apply specific formatting for Source line
        if line.startswith('Time Period of analysis:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans")
            
        # Apply specific formatting for Source line
        if line.startswith('News search:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=4, start_column=1, end_row=4, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans")
            

def add_styling_to_worksheet(ws, df, start_row, comment):
    # Apply table heading as comment
    cell = ws.cell(row=start_row, column=1)
    cell.value = comment
    cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
    cell.font = Font(color="000000", bold=True, name="Gill Sans")
    cell.alignment = Alignment(horizontal='center')
    ws.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=len(df.columns))
    
    # Increment the start row
    start_row += 1

    # Apply styling to column headers
    for col_idx, col_name in enumerate(df.columns, start=1):
        cell = ws.cell(row=start_row, column=col_idx)
        cell.value = col_name
        cell.font = Font(color="000000", bold=True ,name="Gill Sans")
        cell.alignment = Alignment(horizontal='center')
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))  
        
    start_row += 1

    # Write DataFrame values with styling
    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=False), start=start_row):
        for c_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx)
            if isinstance(value, pd.Period):
                cell.value = value.strftime('%Y-%m') 
            else:
                cell.value = value
            cell.font = Font(name="Gill Sans")    
            cell.alignment = Alignment(horizontal='center')
    
    # Apply borders to all cells
    for row in ws.iter_rows(min_row=start_row, max_row=start_row+len(df), min_col=1, max_col=len(df.columns)):
        for cell in row:
            cell.border = Border(left=Side(border_style="thin", color="000000"),
                                 right=Side(border_style="thin", color="000000"),
                                 top=Side(border_style="thin", color="000000"),
                                 bottom=Side(border_style="thin", color="000000"))
            
def multiple_dfs(df_list, sheet_name, file_name, comments, entity_info):
    wb = Workbook()
    ws = wb.active
    current_row = 1
    
    # Add entity information to the first 4 rows
    add_entity_info(ws, entity_info, current_row)
    current_row += 6
    
    for df, comment in zip(df_list, comments):
        add_styling_to_worksheet(ws, df, current_row, comment)
        current_row += len(df) + 4
    
    wb.save(file_name)


def add_table_to_slide(slide, df, title, textbox_text):
    rows, cols = df.shape
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(14)
    max_table_height = Inches(7)
    total_height_needed = Inches(0.8 * (rows + 1))
    height = max_table_height if total_height_needed > max_table_height else total_height_needed

    # Add title shape (above the table)
    title_shape = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.2))
    title_frame = title_shape.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add the table
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = df.columns[i]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Gill Sans'
                run.font.size = Pt(17)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.values[i, j])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Gill Sans'
                    run.font.size = Pt(15)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Add a text box above the table (shared across all DataFrame slides)
    textbox_left = Inches(1)  # Adjust left positioning as needed
    textbox_top = Inches(0.8)  # Adjust top positioning as needed
    textbox_width = Inches(14)  # Adjust width
    textbox_height = Inches(1)  # Adjust height

    text_box = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = textbox_text  # The custom text box content for each slide
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)  # Adjust the font size as needed
#             run.font.bold = True
            run.font.name = 'Gill Sans'
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text

    # Add the image (footer logo) at the bottom of the slide
    left = Inches(0.0)
    top = prs.slide_height - Inches(1)
    # slide.shapes.add_picture( left, top, height=Inches(1))  # Adjust as needed




# # Function to save multiple DataFrames in a single Excel sheet
# def multiple_dfs(df_list, sheets, file_name, spaces, comments):
#     writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
#     row = 2
#     for dataframe, comment in zip(df_list, comments):
#         pd.Series(comment).to_excel(writer, sheet_name=sheets, startrow=row,
#                                     startcol=1, index=False, header=False)
#         dataframe.to_excel(writer, sheet_name=sheets, startrow=row + 1, startcol=0)
#         row = row + len(dataframe.index) + spaces + 2
#     writer.close()
    
def top_10_dfs(df_list, file_name, comments, top_11_flags):
    writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
    row = 2
    for dataframe, comment, top_11_flag in zip(df_list, comments, top_11_flags):
        if top_11_flag:
            top_df = dataframe.head(50)  # Select the top 11 rows for specific DataFrames
        else:
            top_df = dataframe  # Leave other DataFrames unchanged

        top_df.to_excel(writer, sheet_name="Top 10 Data", startrow=row, index=True)
        row += len(top_df) + 2  # Move the starting row down by len(top_df) + 2 rows

    # Create a "Report" sheet with all the DataFrames
    for dataframe, comment in zip(df_list, comments):
        dataframe.to_excel(writer, sheet_name="Report", startrow=row, index=True, header=True)
        row += len(dataframe) + 2  # Move the starting row down by len(dataframe) + 2 rows

    writer.close()    
    

# Streamlit app with a sidebar layout
st.set_page_config(layout="wide")

# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("Meltwater Data Insights Dashboard")

# Sidebar for file upload and download options
st.sidebar.title("Upload a file for tables")

# File Upload Section
file = st.sidebar.file_uploader("Upload Data File (Excel or CSV)", type=["xlsx", "csv"])

if file:
    st.sidebar.write("File Uploaded Successfully!")

    # Load data
    data = load_data(file)

    if data is not None:
        # Data Preview Section (optional)
        # st.write("## Data Preview")
        # st.write(data)

        # Data preprocessing
        data.drop(columns=data.columns[10:], axis=1, inplace=True)
        data = data.rename({'Influencer': 'Journalist'}, axis=1)
        data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
        data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        finaldata = data
        finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()

        # Share of Voice (SOV) Calculation
        En_sov = pd.crosstab(finaldata['Entity'], columns='News Count', values=finaldata['Entity'], aggfunc='count').round(0)
        En_sov.sort_values('News Count', ascending=False)
        En_sov['% '] = ((En_sov['News Count'] / En_sov['News Count'].sum()) * 100).round(2)
        Sov_table = En_sov.sort_values(by='News Count', ascending=False)
        Sov_table.loc['Total'] = Sov_table.sum(numeric_only=True, axis=0)
        Entity_SOV1 = Sov_table.round()
        Entity_SOV3 = pd.DataFrame(Entity_SOV1.to_records()).round()
        Entity_SOV3['% '] = Entity_SOV3['% '].astype(int)
        Entity_SOV3['% '] = Entity_SOV3['% '].astype(str) + '%'

        # Additional MOM DataFrames
        finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()
        sov_dt = pd.crosstab((finaldata['Date'].dt.to_period('M')), finaldata['Entity'], margins=True, margins_name='Total')
        sov_dt1 = pd.DataFrame(sov_dt.to_records())
        
        
        #Publication Name
        pub_table = pd.crosstab(finaldata['Publication Name'], finaldata['Entity'])
        pub_table['Total'] = pub_table.sum(axis=1)
        pubs_table = pub_table.sort_values('Total', ascending=False).round()
        pubs_table.loc['GrandTotal'] = pubs_table.sum(numeric_only=True, axis=0)
        pubs_table = pd.DataFrame(pubs_table.to_records())
        

        PP = pd.crosstab(finaldata['Publication Name'], finaldata['Publication Type'])
        PP['Total'] = PP.sum(axis=1)
        PP_table = PP.sort_values('Total', ascending=False).round()
        PP_table.loc['GrandTotal'] = PP_table.sum(numeric_only=True, axis=0)
        
        #Publication Name & Entity Table
        PT_Entity = pd.crosstab(finaldata['Publication Type'], finaldata['Entity'])
        PT_Entity['Total'] = PT_Entity.sum(axis=1)
        PType_Entity = PT_Entity.sort_values('Total', ascending=False).round()
        PType_Entity.loc['GrandTotal'] = PType_Entity.sum(numeric_only=True, axis=0)
        PType_Entity = pd.DataFrame(PType_Entity.to_records())

        # Journalist Table
        finaldata['Journalist'] = finaldata['Journalist'].str.split(',')
        finaldata = finaldata.explode('Journalist')
        jr_tab = pd.crosstab(finaldata['Journalist'], finaldata['Entity'])
        jr_tab = jr_tab.reset_index(level=0)
        newdata = finaldata[['Journalist', 'Publication Name']]
        Journalist_Table = pd.merge(jr_tab, newdata, how='inner', left_on=['Journalist'], right_on=['Journalist'])
        Journalist_Table.drop_duplicates(subset=['Journalist'], keep='first', inplace=True)
        valid_columns = Journalist_Table.select_dtypes(include='number').columns
        Journalist_Table['Total'] = Journalist_Table[valid_columns].sum(axis=1)
        Jour_table = Journalist_Table.sort_values('Total', ascending=False).round()
        bn_row = Jour_table.loc[Jour_table['Journalist'] == 'Bureau News']
        Jour_table = Jour_table[Jour_table['Journalist'] != 'Bureau News']
        Jour_table = pd.concat([Jour_table, bn_row], ignore_index=True)
#         Jour_table = Journalist_Table.reset_index(drop=True)
        Jour_table.loc['GrandTotal'] = Jour_table.sum(numeric_only=True, axis=0)
        columns_to_convert = Jour_table.columns.difference(['Journalist', 'Publication Name'])
        Jour_table[columns_to_convert] = Jour_table[columns_to_convert].astype(int)
        Jour_table.insert(1, 'Publication Name', Jour_table.pop('Publication Name'))
        
        # Remove square brackets and single quotes from the 'Journalist' column
        data['Journalist'] = data['Journalist'].str.replace(r"^\['(.+)'\]$", r"\1", regex=True)
        # Fill missing values in 'Influencer' column with 'Bureau News'
        data['Journalist'] = data['Journalist'].fillna('Bureau News')

        # Function to classify news exclusivity and topic
        def classify_exclusivity(row):
            entity_name = row['Entity']
            if entity_name.lower() in row['Headline'].lower():
                return "Exclusive"
            else:
                return "Not Exclusive"

        finaldata['Exclusivity'] = finaldata.apply(classify_exclusivity, axis=1)

        # Define a dictionary of keywords for each entity
        entity_keywords = {
            'Amazon': ['Amazon', 'Amazons', 'amazon'],
            # Add other entities and keywords here
        }

        def qualify_entity(row):
            entity_name = row['Entity']
            text = row['Headline']
            if entity_name in entity_keywords:
                keywords = entity_keywords[entity_name]
                if any(keyword in text for keyword in keywords):
                    return "Qualified"
            return "Not Qualified"

        finaldata['Qualification'] = finaldata.apply(qualify_entity, axis=1)

        # Topic classification
        topic_mapping = {
            'Merger': ['merger', 'merges'],
            'Acquire': ['acquire', 'acquisition', 'acquires'],
            'Partnership': ['partnership', 'tie-up'],
            'Business Strategy': ['launch', 'campaign', 'IPO', 'sales'],
            'Investment and Funding': ['invest', 'funding'],
            'Employee Engagement': ['layoff', 'hiring'],
            'Financial Performance': ['profit', 'loss', 'revenue'],
            'Business Expansion': ['expansion', 'opens'],
            'Leadership': ['ceo'],
            'Stock Related': ['stock', 'shares'],
            'Awards & Recognition': ['award'],
            'Legal & Regulatory': ['penalty', 'scam'],
        }

        def classify_topic(headline):
            for topic, words in topic_mapping.items():
                if any(word in headline.lower() for word in words):
                    return topic
            return 'Other'

        finaldata['Topic'] = finaldata['Headline'].apply(classify_topic)

        dfs = [Entity_SOV3, sov_dt1, pubs_table, Jour_table, PType_Entity, PP_table]
        comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table','PubType Entity Table', 'Pub Type and Pub Name Table']

        # Sidebar for download options
        st.sidebar.write("## Download Options")
        download_formats = st.sidebar.selectbox("Select format:", ["Excel", "CSV", "Excel (Entity Sheets)"])

        if st.sidebar.button("Download Data"):
            if download_formats == "Excel":
                # Download all DataFrames as a single Excel file
                excel_io = io.BytesIO()
                with pd.ExcelWriter(excel_io, engine="xlsxwriter") as writer:
                    for df, comment in zip(dfs, comments):
                        df.to_excel(writer, sheet_name=comment, index=False)
                excel_io.seek(0)
                b64_data = base64.b64encode(excel_io.read()).decode()
                href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="data.xlsx">Download Excel</a>'
                st.sidebar.markdown(href_data, unsafe_allow_html=True)

            elif download_formats == "CSV":
                # Download all DataFrames as CSV
                csv_io = io.StringIO()
                for df in dfs:
                    df.to_csv(csv_io, index=False)
                csv_io.seek(0)
                b64_data = base64.b64encode(csv_io.read().encode()).decode()
                href_data = f'<a href="data:text/csv;base64,{b64_data}" download="data.csv">Download CSV</a>'
                st.sidebar.markdown(href_data, unsafe_allow_html=True)

            elif download_formats == "Excel (Entity Sheets)":
                # Download DataFrames as Excel with separate sheets by entity
                excel_io = io.BytesIO()
                with pd.ExcelWriter(excel_io, engine="xlsxwriter") as writer:
                    create_entity_sheets(finaldata, writer)
                excel_io.seek(0)
                b64_data = base64.b64encode(excel_io.read()).decode()
                href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="entity_sheets.xlsx">Download Entity Sheets</a>'
                st.sidebar.markdown(href_data, unsafe_allow_html=True)
                
         
        # Download selected DataFrame
        st.sidebar.write("## Download Selected DataFrame")
        
        dataframes_to_download = {
            "Entity_SOV1": Entity_SOV3,
            "Data": data,
            "Finaldata": finaldata,
            "Month-on-Month":sov_dt1,
            "Publication Table":pubs_table,
            "Journalist Table":Jour_table,
            "Publication Type and Name Table":PP_table,
            "Publication Type Table with Entity":PType_Entity,
            # "Publication type,Publication Name and Entity Table":ppe1,
            "Entity-wise Sheets": finaldata  # Add this option to download entity-wise sheets
        }
        selected_dataframe = st.sidebar.selectbox("Select DataFrame:", list(dataframes_to_download.keys()))
        
        if st.sidebar.button("Download Selected DataFrame"):
            if selected_dataframe in dataframes_to_download:
                # Create a link to download the selected DataFrame in Excel
                selected_df = dataframes_to_download[selected_dataframe]
                excel_io_selected = io.BytesIO()
                with pd.ExcelWriter(excel_io_selected, engine="xlsxwriter", mode="xlsx") as writer:
                    selected_df.to_excel(writer, index=True)
                excel_io_selected.seek(0)
                b64_selected = base64.b64encode(excel_io_selected.read()).decode()
                href_selected = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_selected}" download="{selected_dataframe}.xlsx">Download {selected_dataframe} Excel</a>'
                st.sidebar.markdown(href_selected, unsafe_allow_html=True)
                
                
        # Download All DataFrames as a Single Excel Sheet
        st.sidebar.write("## Download All DataFrames as a Single Excel Sheet")
        file_name_all = st.sidebar.text_input("Enter file name for all DataFrames", "all_dataframes.xlsx")
        
        if st.sidebar.button("Download All DataFrames"):
            # List of DataFrames to save
            dfs = [Entity_SOV3, sov_dt1, pubs_table, Jour_table, PType_Entity, PP_table]
            comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
                        'Pub Type and Entity Table', 'Pub Type and Pub Name Table'
                        ]
            
            entity_info = """Entity:
Time Period of analysis: 19th April 2023 to 18th April 2024
Source: (Online) Meltwater, Select 100 online publications, which include General mainlines, Business and Financial publications, news age media, technology publications.
News search: All Articles: entity mentioned at least once in the article"""
            excel_io_all = io.BytesIO()
            multiple_dfs(dfs, 'Tables', excel_io_all, comments, entity_info)
            excel_io_all.seek(0)
            b64_all = base64.b64encode(excel_io_all.read()).decode()
            href_all = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_all}" download="{file_name_all}">Download All DataFrames Excel</a>'
            st.sidebar.markdown(href_all, unsafe_allow_html=True)

    # Load the image files
    # img_path = r"D:\Akshay.Annaldasula\OneDrive - Adfactors PR Pvt Ltd\Documents\NewLogo.PNG"
    # img_path1 = r"D:\Akshay.Annaldasula\OneDrive - Adfactors PR Pvt Ltd\Pictures\Picture1.png"

    # Create a new PowerPoint presentation with widescreen dimensions
    prs = Presentation()               
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add the first slide with the image
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    # slide.shapes.add_picture( Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Add the text box above the image
    textbox_left = Inches(0.5)  # Adjust the left position as needed
    textbox_top = Inches(5)   # Adjust the top position as needed
    textbox_width = Inches(15)  # Adjust the width as needed
    textbox_height = Inches(1)  # Adjust the height as needed

    text_box = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = "Client Name"
        
    # Set font size to 30 and make the text bold and white
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(50)
            run.font.bold = True
#           run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(255, 255, 255)  # White color
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0.0)  # Adjust the left position as needed
    top = prs.slide_height - Inches(1)  # Adjust the top position as needed
    # slide.shapes.add_picture( left, top, height=Inches(1))  # Adjust the height as needed 

        
    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames

    # Set title text and format for Parameters slide
    header_text = "Parameters"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.7))
    header_frame = header_shape.text_frame
    header_frame.text = header_text

    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    # Add Time Period text
    time_period_text = "Time Period : 19th April 2023 to 18th April 2024"
    time_period_shape = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(14), Inches(0.5))
    time_period_frame = time_period_shape.text_frame
    time_period_frame.text = time_period_text
    # time_period_frame.paragraphs[0].font.bold = True
    time_period_frame.paragraphs[0].font.size = Pt(24)
    time_period_frame.paragraphs[0].font.name = 'Gill Sans'


    # Add Source text
    source_text = "Source : (Online)Meltwater, Select 100 online publications, which include General mainlines, Business and Financial publications, news age media, technology publications."
    source_shape = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))  # Adjusted width
    source_frame = source_shape.text_frame
    source_frame.word_wrap = True  # Enable text wrapping
    p = source_frame.add_paragraph()  # Create a paragraph for text
    p.text = source_text  # Set the text

    p.font.size = Pt(24)
    p.font.name = 'Gill Sans'  # Changed to Arial for compatibility

    # Add News Search text
    news_search_text = "News Search : All Articles: entity mentioned at least once in the article "
    news_search_shape = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(10), Inches(0.75))  # Adjusted width and height
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True  # Enable text wrapping
    p2 = news_search_frame.add_paragraph()  # Create a paragraph for text
    p2.text = news_search_text  # Set the text

    # Set font properties after text is added
    # p2.font.bold = True
    p2.font.size = Pt(24)
    p2.font.name = 'Gill Sans'  # Changed to Arial for compatibility
        
    # Add the first slide with the image
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    # slide.shapes.add_picture( Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Add the text box above the image
    textbox_left = Inches(0.5)  # Adjust the left position as needed
    textbox_top = Inches(5)   # Adjust the top position as needed
    textbox_width = Inches(15)  # Adjust the width as needed
    textbox_height = Inches(1)  # Adjust the height as needed

    text_box = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = "Online Media"

    # Set font size to 30 and make the text bold and white
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(50)
            run.font.bold = True
#           run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(255, 255, 255)  # White color
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0.0)  # Adjust the left position as needed
    top = prs.slide_height - Inches(1)  # Adjust the top position as needed
    # slide.shapes.add_picture( left, top, height=Inches(1))  # Adjust the height as needed 
         
    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames

    # Set title text and format for Parameters slide
    header_text = "Inferences and Recommendations"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(14), Inches(0.7))
    header_frame = header_shape.text_frame
    header_frame.text = header_text
    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  


    # Add SOV text
    sov_text = ("Share of Voice :")
    sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
    sov_text_frame = sov_text_shape.text_frame
    sov_text_frame.word_wrap = True
    sov_text_frame.clear()  # Clear any default paragraph

    p = sov_text_frame.add_paragraph()
    p.text = "Share of Voice :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True

    sov_text = (
    "IIT Ropar and its peers collectively received a total of 5017 prominent news mentions online during the specified time period.\n"
    "Among these, IIT Madras dominates the conversation with 35% of the total SOV, indicating significant media coverage and visibility.\n"
    "IIT Delhi follows IIT Madras, capturing 21% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
    "IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 20%, 17%, and 6% of the SOV respectively.\n"
    "IIT Ropar holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e. last in the SOV.\n"
    "Despite ranking lower in terms of SOV, IIT Ropar's presence indicates some level of visibility and recognition within the online media landscape.\n"
    "Given the relatively lower SOV compared to peers like IIT Delhi, IIT Madras, and others, there are opportunities for IIT Ropar to enhance its online presence and visibility through strategic communications efforts.\n"
    "IIT Ropar has received 239 all mentions and 44 prominent articles in online media and stands last in both the SOVs.\n"
        )
    sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
    sov_text_frame = sov_text_shape.text_frame
    sov_text_frame.word_wrap = True
    sov_text_frame.clear()  # Clear any default paragraph


    p = sov_text_frame.add_paragraph()
    p.text = (
    "IIT Ropar and its peers collectively received a total of 5017 prominent news mentions online during the specified time period.\n"
    "Among these, IIT Madras dominates the conversation with 35% of the total SOV, indicating significant media coverage and visibility.\n"
    "IIT Delhi follows IIT Madras, capturing 21% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
    "IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 20%, 17%, and 6% of the SOV respectively.\n"
    "IIT Ropar holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e. last in the SOV.\n"
    "Despite ranking lower in terms of SOV, IIT Ropar's presence indicates some level of visibility and recognition within the online media landscape.\n"
    "Given the relatively lower SOV compared to peers like IIT Delhi, IIT Madras, and others, there are opportunities for IIT Ropar to enhance its online presence and visibility through strategic communications efforts.\n"
    "IIT Ropar has received 239 all mentions and 44 prominent articles in online media and stands last in both the SOVs.\n"
    )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'

    # Add Source text
    source_text = ("Publications :")
    source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.8), Inches(14), Inches(1))
    source_frame = source_shape.text_frame
    source_frame.word_wrap = True
    source_frame.clear()  # Clear any default paragraph
    p = source_frame.add_paragraph()
    p.text = "Publications :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True


    source_text = (
    "The leading publications reporting on IIT Ropar and its competitors are Times of India, contributing 561 articles, followed by Economic Times with 467 articles, and The Indian Express with 455 articles.\n"
"Among these ,publications covering news on IIT Ropar specifically are The Indian Express takes the lead with 9 articles, followed by Tribune with 9 articles, and Times of India with 7 articles.\n"
"The top 10 publications writing articles on IIT Ropar contribute 86% of the total 44 articles.\n" 
)
    source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.1), Inches(14), Inches(1))
    source_frame = source_shape.text_frame
    source_frame.word_wrap = True
    source_frame.clear()  # Clear any default paragraph
    p = source_frame.add_paragraph()
    p.text = (
    "The leading publications reporting on IIT Ropar and its competitors are Times of India, contributing 561 articles, followed by Economic Times with 467 articles, and The Indian Express with 455 articles.\n"
"Among these ,publications covering news on IIT Ropar specifically are The Indian Express takes the lead with 9 articles, followed by Tribune with 9 articles, and Times of India with 7 articles.\n"
"The top 10 publications writing articles on IIT Ropar contribute 86% of the total 44 articles.\n" 
)
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'

    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)


    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames


    # Set title text and format for Parameters slide
    header_text = "Inferences and Recommendations"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
    header_frame = header_shape.text_frame
    header_frame.text = header_text 
    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


    # Add News Search text
    news_search_text = ("Journalists :")
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = "Journalists :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True

    # Add News Search text
    news_search_text = ("The top journalists reporting on IIT Ropar and its competitors are Sukanya Nandy from News18 with 59 articles, followed by Hemali Chapia from TOI with 44 articles, and Suramya Sunilraj from News18 with 43 articles.\n"
                    "Among the journalists specifically covering IIT Ropar are Munieshwer A Sagar and Deepak Yadav from Times of Indian  has authored 1 articles each  and Arushi Mishra from Hindu Business Line written 1 article.\n"
                    "IIT Ropar has received a total of 44 articles in news coverage. Among these, 39 i.e 88% of the articles were filed by Bureaus, while the remaining 5 i.e 12% were written by individual journalists.\n"
                    "A total of 387 journalists have written 1155 articles covering IIT Ropar and its competitors.\n"
                    "Out of which, 5 journalists have specifically written 5 articles mentioning IIT Ropar i.e of the total journalists writing on IIT Ropar and its competitors only 1% of them have mentioned IIT Ropar in their articles.\n"
                    "While this constitutes a very less number, there is an large opportunity for IIT Ropar to engage with the remaining 882 journalists to enhance its news coverage and reach.\n"
                   )
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = ("The top journalists reporting on IIT Ropar and its competitors are Sukanya Nandy from News18 with 59 articles, followed by Hemali Chapia from TOI with 44 articles, and Suramya Sunilraj from News18 with 43 articles.\n"
                    "Among the journalists specifically covering IIT Ropar are Munieshwer A Sagar and Deepak Yadav from Times of Indian  has authored 1 articles each  and Arushi Mishra from Hindu Business Line written 1 article.\n"
                    "IIT Ropar has received a total of 44 articles in news coverage. Among these, 39 i.e 88% of the articles were filed by Bureaus, while the remaining 5 i.e 12% were written by individual journalists.\n"
                    "A total of 387 journalists have written 1155 articles covering IIT Ropar and its competitors.\n"
                    "Out of which, 5 journalists have specifically written 5 articles mentioning IIT Ropar i.e of the total journalists writing on IIT Ropar and its competitors only 1% of them have mentioned IIT Ropar in their articles.\n"
                    "While this constitutes a very less number, there is an large opportunity for IIT Ropar to engage with the remaining 882 journalists to enhance its news coverage and reach.\n"
                   )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'

    # Add News Search text
    news_search_text = ("Publication Types :" )
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.6), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = "Publication Type :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True

    news_search_text = ("Top Publication Types writing on IIT Ropar are General and  Business & Financials they both contribute 90% of the total news coverage on IIT Ropar.\n"
"IIT Madras and IIT Delhi dominates across all publication types, especially in general, business ,technology, and digital-first publications.\n"
"IIT Ropar may find value in engaging more with General and Business along with technology, and digital-first publications to expand her reach and visibility among broader audiences.\n"
                   )
    news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.0), Inches(14), Inches(0.5))
    news_search_frame = news_search_shape.text_frame
    news_search_frame.word_wrap = True
    news_search_frame.clear()  # Clear any default paragraph
    p = news_search_frame.add_paragraph()
    p.text = ("Top Publication Types writing on IIT Ropar are General and  Business & Financials they both contribute 90% of the total news coverage on IIT Ropar.\n"
"IIT Madras and IIT Delhi dominates across all publication types, especially in general, business ,technology, and digital-first publications.\n"
"IIT Ropar may find value in engaging more with General and Business along with technology, and digital-first publications to expand her reach and visibility among broader audiences.\n"
                   )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'
        
    # Add title slide after the first slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Clear existing placeholders
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.clear()  # Clear existing text frames
        
    # Set title text and format for Parameters slide
    header_text = "Inferences and Recommendations"
    header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
    header_frame = header_shape.text_frame
    header_frame.text = header_text
    for paragraph in header_frame.paragraphs:
        for run in paragraph.runs:
            run.text = header_text
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
            # Set alignment to center
            paragraph.alignment = PP_ALIGN.CENTER
            # Set vertical alignment to be at the top
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


    # # Add Time Period text
    time_period_text = ("Monthly Coverage , Peak and Topics :")
    time_period_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
    time_period_frame = time_period_shape.text_frame
    time_period_frame.text = time_period_text
    time_period_frame.word_wrap = True
    time_period_frame.clear() 

    p = time_period_frame.add_paragraph()
    p.text = "Monthly Coverage , Peak and Topics :"
    p.font.size = Pt(20)
    p.font.name = 'Gill Sans'
    p.font.underline = True
    p.font.bold = True


    time_period_text = ("The total number of news articles over the period is 5311. There is noticeable fluctuation in coverage from month to month, with periods of both increase and decline.\n"
"Sept 2023 saw the highest number of articles, with 524 mentions. This spike suggests a significant event or increased media focus on Steel Industry during that month.\n"
"Dec 2023 has the lowest coverage so far, with only 337 mentions.\n"
"There was  peak in Sept-23 due to following news:The increase in volume is due to Prime Minister Narendra Modi's visit to Chhattisgarh and the launch of multiple development projects, including the dedication of the NMDC Steel Plant in Bastar. The projects are expected to provide employment opportunities and contribute to the reduction of debt burden. The political controversy surrounding the Nagarnar Steel Plant and the statements made by Chief Minister Bhupesh Baghel have also contributed to the spike in volume\n"
                   )
    time_period_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(14), Inches(0.5))
    time_period_frame = time_period_shape.text_frame
    time_period_frame.text = time_period_text
    time_period_frame.word_wrap = True
    time_period_frame.clear() 

    p = time_period_frame.add_paragraph()
    p.text = ("The total number of news articles over the period is 5311. There is noticeable fluctuation in coverage from month to month, with periods of both increase and decline.\n"
"Sept 2023 saw the highest number of articles, with 524 mentions. This spike suggests a significant event or increased media focus on Steel Industry during that month.\n"
"Dec 2023 has the lowest coverage so far, with only 337 mentions.\n"
"There was  peak in Sept-23 due to following news:The increase in volume is due to Prime Minister Narendra Modi's visit to Chhattisgarh and the launch of multiple development projects, including the dedication of the NMDC Steel Plant in Bastar. The projects are expected to provide employment opportunities and contribute to the reduction of debt burden. The political controversy surrounding the Nagarnar Steel Plant and the statements made by Chief Minister Bhupesh Baghel have also contributed to the spike in volume\n"
                   )
    p.font.size = Pt(18)
    p.font.name = 'Gill Sans'


    # Sidebar for PowerPoint download settings
    st.sidebar.write("## Download All DataFrames as a PowerPoint File")
    pptx_file_name = st.sidebar.text_input("Enter file name for PowerPoint", "dataframes_presentation.pptx")

    if st.sidebar.button("Download PowerPoint"):
        # List of DataFrames to save
        dfs = [Entity_SOV3, sov_dt1, pubs_table, Jour_table, PType_Entity, PP_table]
        table_titles = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
                    'Pub Type and Entity Table', 'Pub Type and Pub Name Table',
                    ]
        textbox_text = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
                    'Pub Type and Entity Table', 'Pub Type and Pub Name Table',
                    ]
        # Create a new PowerPoint presentation
        prs = Presentation()

        # Loop through each DataFrame and create a new slide with a table
        for i, (df, title) in enumerate(zip(dfs, table_titles)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_table_to_slide(slide, df, title, textbox_text[i])

        # Save presentation to BytesIO for download
        pptx_output = io.BytesIO()
        prs.save(pptx_output)
        pptx_output.seek(0)

        # Provide download button
        st.sidebar.download_button(
            label="Download PowerPoint Presentation",
            data=pptx_output,
            file_name=pptx_file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
#         # Download All DataFrames as a Single Excel Sheet
#         st.sidebar.write("## Download All DataFrames as a Single Excel Sheet")
#         file_name_all = st.sidebar.text_input("Enter file name for all DataFrames", "all_dataframes.xlsx")
# #         download_options = st.sidebar.selectbox("Select Download Option:", [ "Complete Dataframes"])
        
#         if st.sidebar.button("Download All DataFrames"):
#             # List of DataFrames to save
#             dfs = [Entity_SOV1, sov_dt, pubs_table, Jour_table, PType_Entity, PP_table, ppe1]
#             comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
#                         'Pub Type and Entity Table', 'Pub Type and Pub Name Table',
#                         'PubType PubName and Entity Table']
            
#             excel_path_all = os.path.join(download_path, file_name_all)
#             multiple_dfs(dfs, 'Tables', excel_path_all, 2, comments)
#             st.sidebar.write(f"All DataFrames saved at {excel_path_all}")

#         # Loop through each dataframe and create a new slide for each one
#         for i, (df, title) in enumerate(zip(dfrs, table_titles)):
#             slide = prs.slides.add_slide(prs.slide_layouts[6])
#             add_table_to_slide(slide, df, title, textbox_text[i])

else:
    st.sidebar.write("No file uploaded yet.")
